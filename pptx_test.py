from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.slide import SlideLayouts

prs = Presentation()
blank_slide_layout = prs.slide_layouts[6]
print(prs.slide_height)
print(prs.slide_width)#10inch

slide = prs.slides.add_slide(blank_slide_layout)
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(0, 0, 0)
# 검은 화면


for i in range(4):
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)
    left = width = height = Inches(1)
    top = Inches(5)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.size = Pt(40)
    tf.text = "This is text inside a textbox"


    p = tf.add_paragraph()
    p.text = "This is a second paragraph that's bold한글?"
    p.font.bold = True

    p = tf.add_paragraph()
    p.text = "This is a third paragraph that's big"
    p.font.size = Pt(40)
testSlide = prs.slides[1]
tbackgroud = testSlide.background
fill = tbackgroud.fill
fill.solid()
fill.fore_color.rgb = RGBColor(0, 0, 0)

prs.save('test.pptx')