from pptx import Presentation
import codecs
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def get_size(prs):
    prs = Presentation("smcc_youth.pptx")
    #914400 EMU / Inches
    print(prs.slide_height)
    print(prs.slide_width)

def create_title(title, left = 0.5, top = 0.2, width = 11.333, height = 1.2):
    # 빈 슬라이드 생성
    empty_slide = prs.slide_layouts[6]
    slide = prs.slides.add_slide(empty_slide)

    # 배경색 지정
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0, 0, 0)#black

    #텍스트 상자 위치, 크기
    tleft = Inches(left)
    ttop = Inches(top)
    twidth = Inches(width)
    theight = Inches(height)

    # 텍스트 박스 생성
    text_box = slide.shapes.add_textbox(tleft, ttop, twidth, theight)
    text_paragraph = text_box.text_frame.paragraphs[0]
    text_paragraph.alignment = PP_ALIGN.LEFT#왼쪽 정렬

    # 텍스트 상자 입력
    tr = text_paragraph.add_run()
    tr.text = title.strip()

    # 글자 스타일 설정
    font = tr.font
    font.size = Pt(45)#font size
    font.color.rgb = RGBColor(255, 255, 255)#font color -> white

def create_lyrics(lyric, left = 0, top = 6, width = 13.333, height = 1.5):
    empty_slide = prs.slide_layouts[6]
    slide = prs.slides.add_slide(empty_slide)

    # 배경 관리(검은색)
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0, 0, 0)#black

    #shape 위치, 크기
    sleft = Inches(left)
    stop = Inches(top)
    swidth = Inches(width)
    sheight = Inches(height)

    # shape 생성
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, sleft, stop, swidth, sheight)
    shape_fill = shape.fill
    shape_fill.solid()
    shape_fill.fore_color.rgb = RGBColor(175, 180, 185)
    shape_border = shape.line
    shape_border.color.rgb = RGBColor(0, 0, 0)

    #텍스트 상자 위치, 크기    
    tleft = Inches(1)
    ttop = Inches(6)
    twidth = Inches(11.333)
    theight = Inches(1.4)

    # 텍스트 상자 생성
    text_box = slide.shapes.add_textbox(tleft, ttop, twidth, theight)
    tf = text_box.text_frame
    text_p = tf.paragraphs[0]        
    text_p.alignment = PP_ALIGN.CENTER
    
    tr = text_p.add_run()
    tr.text = lyric

    # 가사 스타일 설정
    font = tr.font    
    font.size = Pt(44)
    font.bold = True
    # font.name = "Nanum JangMiCe"
    font.name = "나눔손글씨 장미체"
    font.color.rgb = RGBColor(255, 255, 255)#white

###################################__main__###############################################

# create pptx
prs = Presentation("basic.pptx")
text_path = "D:/python/수련회찬양.txt"
f = codecs.open(text_path, "r", "utf-8")
lines = f.readlines()

title = ""
lyric = ""
is_title_slide = False

for line in lines:
    if line.startswith("##"):#타이틀 생성
        title = line[1:]
        is_title_slide = True
    elif line.startswith("#"):#가사 슬라이드 생성
        if is_title_slide:
            create_title(title)
            title = ""
            is_title_slide = False
        else:
            create_lyrics(lyric)
            lyric = ""
    else:#가사 이어 붙이기
        if len(lyric) > 4:
            lyric = lyric +"\n"+line.strip()
        else:
            lyric = lyric + line.strip()

prs.save("smcc_youth.pptx")